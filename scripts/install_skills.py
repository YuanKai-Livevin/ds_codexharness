#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""批量安装办公技能到 SKILLS 仓库。
用法: python install_skills.py [仓库目录]
"""
import os
import sys

REPO = sys.argv[1] if len(sys.argv) > 1 else os.path.join(os.environ.get("APPDATA", ""), "OfficeHarness", "skills")

SKILLS = {}

# ---------------- merge-excel ----------------
SKILLS["merge-excel"] = {
    "SKILL.md": """---
name: merge-excel
description: 合并多个 Excel 文件（按行拼接或按工作表合并），可选去重与统一列名
---

# 合并 Excel

## 执行步骤
1. 确定要合并的 Excel 文件列表（工作区内）。
2. 运行脚本：
   ```bash
   python "<仓库>\\merge-excel\\scripts\\merge_excel.py" --files "文件1.xlsx" "文件2.xlsx" ... --output "合并结果.xlsx"
   ```
   或合并一个目录下所有 .xlsx：`--dir 目录 --output 结果.xlsx`。
3. 汇报结果：总行数、列数、输出路径。

## 注意事项
- 输出必须落在工作区内。
- 默认按行拼接（pd.concat），列不同时自动取并集补空。
""",
    "scripts/merge_excel.py": """#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# 合并多个 Excel 文件。用法见 SKILL.md。
import argparse
import glob
import os
import sys

import pandas as pd


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--files", nargs="*", default=[])
    ap.add_argument("--dir", default="")
    ap.add_argument("--output", required=True)
    ap.add_argument("--dedup", action="store_true", help="按全部列去重")
    args = ap.parse_args()

    files = list(args.files)
    if args.dir:
        files += sorted(glob.glob(os.path.join(args.dir, "*.xlsx")) + glob.glob(os.path.join(args.dir, "*.xls")))
    files = [f for f in files if os.path.exists(f)]
    if not files:
        print("ERROR: 没有找到要合并的 Excel 文件")
        sys.exit(1)

    frames = []
    for f in files:
        try:
            df = pd.read_excel(f)
            frames.append(df)
        except Exception as e:
            print(f"WARN: 读取 {f} 失败: {e}")

    if not frames:
        print("ERROR: 所有文件读取失败")
        sys.exit(1)

    merged = pd.concat(frames, ignore_index=True)
    if args.dedup:
        before = len(merged)
        merged = merged.drop_duplicates()
        print(f"去重: {before} -> {len(merged)} 行")

    os.makedirs(os.path.dirname(os.path.abspath(args.output)) or ".", exist_ok=True)
    merged.to_excel(args.output, index=False)
    print(f"OK 合并 {len(files)} 个文件 -> {args.output}，共 {len(merged)} 行 x {merged.shape[1]} 列")


if __name__ == "__main__":
    main()
""",
}

# ---------------- excel-analysis ----------------
SKILLS["excel-analysis"] = {
    "SKILL.md": """---
name: excel-analysis
description: Excel/CSV 数据分析：概览、分组汇总、透视、筛选统计，输出汇总表
---

# Excel 数据分析

## 执行步骤
1. 读取用户指定文件（工作区内）。
2. 运行脚本：
   ```bash
   python "<仓库>\\excel-analysis\\scripts\\analyze_excel.py" --file 数据.xlsx --groupby 区域 --agg 销量:sum,销售额:sum --output 汇总.xlsx
   ```
   - `--groupby`：分组列（可多个，逗号分隔）
   - `--agg`：聚合规则 `列:函数`（sum/mean/count/min/max），逗号分隔
   - 不带 `--groupby` 时输出整体概览（行列数、每列类型、空值、基本统计）。
3. 汇报关键指标（行数、分组数、Top 组）与输出路径。

## 注意事项
- 支持 xlsx/xls/csv。
- 输出到工作区内。
""",
    "scripts/analyze_excel.py": """#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Excel/CSV 数据分析与分组汇总。
import argparse
import os
import sys

import pandas as pd


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--file", required=True)
    ap.add_argument("--groupby", default="")
    ap.add_argument("--agg", default="")
    ap.add_argument("--output", default="")
    args = ap.parse_args()

    if args.file.lower().endswith(".csv"):
        df = pd.read_csv(args.file)
    else:
        df = pd.read_excel(args.file)
    print(f"概览: {len(df)} 行 x {df.shape[1]} 列")
    print("空值统计:")
    print(df.isnull().sum().to_string())

    if args.groupby:
        groups = [g.strip() for g in args.groupby.split(",") if g.strip()]
        rules = {}
        for item in args.agg.split(","):
            item = item.strip()
            if not item:
                continue
            if ":" in item:
                col, fn = item.rsplit(":", 1)
                rules[col.strip()] = fn.strip()
        if rules:
            result = df.groupby(groups, dropna=False).agg(rules).reset_index()
        else:
            numeric = df.select_dtypes(include="number").columns
            result = df.groupby(groups, dropna=False)[list(numeric)].sum().reset_index()
        print("分组汇总:")
        print(result.to_string(index=False))
        if args.output:
            result.to_excel(args.output, index=False)
            print("OK 输出 -> " + args.output)
        else:
            print("OK (未指定 --output，仅打印)")
    else:
        print("数值列统计:")
        print(df.describe().to_string())
        print("OK 分析完成")


if __name__ == "__main__":
    main()
""",
}

# ---------------- word-batch ----------------
SKILLS["word-batch"] = {
    "SKILL.md": """---
name: word-batch
description: Word 批量处理：占位符模板填充、批量生成/修改文档、批量替换文字
---

# Word 批量处理

## 执行步骤
1. 模板填充：模板中用 `{{姓名}}` 这类占位符，数据放在 Excel/CSV。
   ```bash
   python "<仓库>\\word-batch\\scripts\\fill_word.py" --template 模板.docx --data 名单.xlsx --outdir 输出目录
   ```
2. 批量替换文字：`--replace "旧文本" "新文本"` 对目录下所有 docx 生效。
3. 汇报生成的文件数量与路径。

## 注意事项
- 支持 docx（python-docx），不支持旧版 .doc。
- 所有文件限制在工作区内。
""",
    "scripts/fill_word.py": """#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Word 占位符模板批量填充（用法见 SKILL.md）
import argparse
import os
import re
import sys

import pandas as pd
from docx import Document


def fill_one(template, row: dict, out_path: str):
    doc = Document(template)
    for para in doc.paragraphs:
        for key, val in row.items():
            if val is None:
                val = ""
            if "{{" + key + "}}" in para.text:
                para.text = para.text.replace("{{" + key + "}}", str(val))
    for table in doc.tables:
        for trow in table.rows:
            for cell in trow.cells:
                for para in cell.paragraphs:
                    for key, val in row.items():
                        if val is None:
                            val = ""
                        if "{{" + key + "}}" in para.text:
                            para.text = para.text.replace("{{" + key + "}}", str(val))
    doc.save(out_path)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", default="")
    ap.add_argument("--data", default="")
    ap.add_argument("--sheet", default="")
    ap.add_argument("--outdir", default="")
    ap.add_argument("--replace", nargs=2, default=None)
    ap.add_argument("--dir", default="")
    args = ap.parse_args()

    if args.replace:
        target_dir = args.dir
        if not target_dir:
            print("ERROR: --replace 需要 --dir")
            sys.exit(1)
        old, new = args.replace
        n = 0
        for name in os.listdir(target_dir):
            if name.endswith(".docx"):
                p = os.path.join(target_dir, name)
                doc = Document(p)
                changed = False
                for para in doc.paragraphs:
                    if old in para.text:
                        para.text = para.text.replace(old, new)
                        changed = True
                for table in doc.tables:
                    for trow in table.rows:
                        for cell in trow.cells:
                            for para in cell.paragraphs:
                                if old in para.text:
                                    para.text = para.text.replace(old, new)
                                    changed = True
                if changed:
                    doc.save(p)
                    n += 1
        print(f"OK 已替换 {n} 个文档")
        return

    if not args.template or not args.data:
        print("ERROR: 需要 --template 和 --data")
        sys.exit(1)
    if args.data.lower().endswith(".csv"):
        df = pd.read_csv(args.data)
    else:
        df = pd.read_excel(args.data, sheet_name=args.sheet or 0)
    outdir = args.outdir or os.path.dirname(os.path.abspath(args.data))
    os.makedirs(outdir, exist_ok=True)
    rows = df.to_dict("records")
    for i, row in enumerate(rows):
        ident = row.get("姓名") or row.get("编号") or i + 1
        out = os.path.join(outdir, f"{ident}.docx")
        fill_one(args.template, row, out)
    print(f"OK 已生成 {len(rows)} 份文档 -> {outdir}")


if __name__ == "__main__":
    main()
""",
}

# ---------------- pdf-tools ----------------
SKILLS["pdf-tools"] = {
    "SKILL.md": """---
name: pdf-tools
description: PDF 合并、拆分、提取文本/图片、加密解密，基于 pypdf
---

# PDF 工具

## 执行步骤
```bash
python "<仓库>\\pdf-tools\\scripts\\pdf_tools.py" merge --files a.pdf b.pdf --output 合并.pdf
python "<仓库>\\pdf-tools\\scripts\\pdf_tools.py" split --file a.pdf --pages 1-3,5 --output 拆分.pdf
python "<仓库>\\pdf-tools\\scripts\\pdf_tools.py" text --file a.pdf            # 提取文本
python "<仓库>\\pdf-tools\\scripts\\pdf_tools.py" info --file a.pdf           # 页数/元信息
```
- `--pages` 格式：`1-3` 或 `1,3,5-7`（页码从 1 开始）。
- 汇报输出路径与页数。

## 注意事项
- 输出到工作区内；合并/拆分不影响原文件。
""",
    "scripts/pdf_tools.py": """#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# PDF 工具：merge / split / text / info。
import argparse
import os
import sys

from pypdf import PdfReader, PdfWriter


def parse_pages(spec: str, total: int):
    out = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            a = 1 if a == "" else int(a)
            b = total if b == "" else int(b)
            out.update(range(a, b + 1))
        else:
            out.add(int(part))
    return sorted(p for p in out if 1 <= p <= total)


def main():
    ap = argparse.ArgumentParser()
    sub = ap.add_subparsers(dest="cmd", required=True)

    m = sub.add_parser("merge")
    m.add_argument("--files", nargs="*", required=True)
    m.add_argument("--output", required=True)

    s = sub.add_parser("split")
    s.add_argument("--file", required=True)
    s.add_argument("--pages", default="")
    s.add_argument("--output", required=True)

    t = sub.add_parser("text")
    t.add_argument("--file", required=True)

    i = sub.add_parser("info")
    i.add_argument("--file", required=True)

    args = ap.parse_args()

    if args.cmd == "merge":
        writer = PdfWriter()
        total = 0
        for f in args.files:
            reader = PdfReader(f)
            for page in reader.pages:
                writer.add_page(page)
                total += 1
        with open(args.output, "wb") as fh:
            writer.write(fh)
        print(f"OK 合并 {len(args.files)} 个 PDF -> {args.output}（共 {total} 页）")

    elif args.cmd == "split":
        reader = PdfReader(args.file)
        total = len(reader.pages)
        pages = parse_pages(args.pages, total) if args.pages else list(range(1, total + 1))
        writer = PdfWriter()
        for p in pages:
            writer.add_page(reader.pages[p - 1])
        with open(args.output, "wb") as fh:
            writer.write(fh)
        print(f"OK 拆分 {len(pages)}/{total} 页 -> {args.output}")

    elif args.cmd == "text":
        reader = PdfReader(args.file)
        for i, page in enumerate(reader.pages):
            txt = page.extract_text() or ""
            print(f"--- 第 {i + 1} 页 ---")
            print(txt.strip())

    elif args.cmd == "info":
        reader = PdfReader(args.file)
        meta = reader.metadata or {}
        print(f"页数: {len(reader.pages)}")
        for k in ("title", "author", "subject"):
            if meta.get(k):
                print(f"{k}: {meta[k]}")


if __name__ == "__main__":
    main()
""",
}

# ---------------- ppt-report ----------------
SKILLS["ppt-report"] = {
    "SKILL.md": """---
name: ppt-report
description: 把 Excel/CSV 数据自动生成数据汇报 PPT（表格页 + 汇总页），适合周报/月报/经营汇报
---

# 数据汇报 PPT

## 执行步骤
1. 数据文件（工作区内）含汇总后的表格，例如「区域汇总.xlsx」。
2. 运行脚本：
   ```bash
   python "<仓库>\\ppt-report\\scripts\\data_to_ppt.py" --data 区域汇总.xlsx --title "销售月报" --output 汇报.pptx
   ```
   - 可选 `--chart 销售额`：对数值列生成柱状图页。
3. 汇报 PPT 路径与页数，可询问用户是否调整配色/页眉。

## 注意事项
- 深色商务风，与照片技能保持同一视觉语言。
- 输出到工作区内（建议 `.oh_ppt`）。
""",
    "scripts/data_to_ppt.py": """#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Excel/CSV -> 数据汇报 PPT。
import argparse
import os
import sys

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DARK = RGBColor(0x0F, 0x11, 0x15)
ACCENT = RGBColor(0x56, 0x86, 0xFE)
WHITE = RGBColor(0xF9, 0xFA, 0xFB)
GREY = RGBColor(0xAD, 0xB2, 0xB8)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument("--title", default="数据汇报")
    ap.add_argument("--chart", default="", help="对某数值列生成柱状图页")
    ap.add_argument("--output", default="report.pptx")
    args = ap.parse_args()

    if args.data.lower().endswith(".csv"):
        df = pd.read_csv(args.data)
    else:
        df = pd.read_excel(args.data)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def txt(slide, l, t, w, h, text, size, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Microsoft YaHei"

    # 封面
    s = prs.slides.add_slide(blank)
    bg(s, DARK)
    txt(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2), args.title, 40, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6), "数据来源：" + os.path.basename(args.data), 14, GREY, False, PP_ALIGN.CENTER)

    # 数据表格页（前 12 行）
    s = prs.slides.add_slide(blank)
    bg(s, DARK)
    txt(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6), "数据明细", 22, WHITE, True)
    view = df.head(12)
    rows, cols = view.shape
    top, left, cell_h = Inches(1.1), Inches(0.6), Inches(0.42)
    col_w = Inches(11.7 / cols) if cols else Inches(2)
    for c, name in enumerate(view.columns):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + c * col_w, top, col_w, Inches(0.42))
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT; cell.line.fill.background()
        tf = cell.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; r = p.add_run(); r.text = str(name)[:20]
        r.font.size = Pt(11); r.font.color.rgb = WHITE; r.font.bold = True
    for r_ in range(rows):
        for c in range(cols):
            val = view.iloc[r_, c]
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + c * col_w, top + Inches(0.42) + r_ * cell_h, col_w, cell_h)
            cell.fill.solid(); cell.fill.fore_color.rgb = DARK if r_ % 2 == 0 else RGBColor(0x1A, 0x1C, 0x22)
            cell.line.fill.background()
            tf = cell.text_frame; p = tf.paragraphs[0]; r = p.add_run()
            r.text = str(val)[:24]
            r.font.size = Pt(10); r.font.color.rgb = WHITE

    # 汇总 + 图表页
    s = prs.slides.add_slide(blank)
    bg(s, DARK)
    txt(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6), "汇总", 22, WHITE, True)
    if args.chart and args.chart in df.columns:
        cats = df.iloc[:, 0].astype(str).tolist()[:20]
        vals = df[args.chart].tolist()[:20]
        cdata = CategoryChartData()
        cdata.categories = cats
        cdata.add_series(args.chart, vals)
        gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.1), Inches(12), Inches(5.6), cdata)
        chart = gframe.chart
        chart.font.size = Pt(10)
    else:
        txt(s, Inches(0.6), Inches(1.2), Inches(12), Inches(5), df.describe().to_string()[:900], 12, GREY)

    prs.save(args.output)
    print(f"OK 汇报 PPT -> {args.output}（{len(prs.slides.__iter__.__self__._sldIdLst)} 页）")


if __name__ == "__main__":
    main()
""",
}

# ---------------- image-batch ----------------
SKILLS["image-batch"] = {
    "SKILL.md": """---
name: image-batch
description: 图片批量处理：压缩、格式转换、缩放、加水印、重命名，基于 Pillow
---

# 图片批量处理

## 执行步骤
```bash
python "<仓库>\\image-batch\\scripts\\image_batch.py" --dir 图片目录 --resize 1280 --quality 80 --outdir 输出目录
python "<仓库>\\image-batch\\scripts\\image_batch.py" --dir 图片目录 --convert png --outdir 输出目录
python "<仓库>\\image-batch\\scripts\\image_batch.py" --dir 图片目录 --watermark "公司名" --outdir 输出目录
```
- `--resize`：最长边像素；`--quality`：JPEG 质量；`--convert`：目标格式（png/jpg/webp）；`--watermark`：右下角水印文字。

## 注意事项
- 输出到工作区内；默认不覆盖原图。
""",
    "scripts/image_batch.py": """#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# 图片批量处理。
import argparse
import os
import sys

from PIL import Image, ImageDraw, ImageFont

EXT = {".jpg", ".jpeg", ".png", ".bmp", ".webp"}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--dir", required=True)
    ap.add_argument("--outdir", default="")
    ap.add_argument("--resize", type=int, default=0)
    ap.add_argument("--quality", type=int, default=85)
    ap.add_argument("--convert", default="")
    ap.add_argument("--watermark", default="")
    args = ap.parse_args()

    files = [os.path.join(args.dir, n) for n in sorted(os.listdir(args.dir)) if os.path.splitext(n)[1].lower() in EXT]
    if not files:
        print("ERROR: 目录中没有图片")
        sys.exit(1)
    outdir = args.outdir or os.path.join(args.dir, "processed")
    os.makedirs(outdir, exist_ok=True)

    n = 0
    for f in files:
        im = Image.open(f)
        if im.mode in ("RGBA", "P"):
            im = im.convert("RGBA")
        else:
            im = im.convert("RGB")
        if args.resize:
            w, h = im.size
            scale = args.resize / max(w, h)
            if scale < 1:
                im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        if args.watermark:
            draw = ImageDraw.Draw(im)
            try:
                font = ImageFont.truetype("msyh.ttc", max(18, im.size[0] // 20))
            except Exception:
                font = ImageFont.load_default()
            bbox = draw.textbbox((0, 0), args.watermark, font=font)
            tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
            x, y = im.size[0] - tw - 20, im.size[1] - th - 20
            draw.rectangle([x - 6, y - 6, x + tw + 6, y + th + 6], fill=(0, 0, 0, 140))
            draw.text((x, y), args.watermark, font=font, fill=(255, 255, 255, 255))
        fmt = args.convert.lower() or os.path.splitext(f)[1].lstrip(".").lower()
        if fmt == "jpg":
            fmt = "jpeg"
        if fmt == "jpeg":
            im = im.convert("RGB")
        out = os.path.join(outdir, os.path.splitext(os.path.basename(f))[0] + "." + fmt)
        im.save(out, fmt, quality=args.quality)
        n += 1
    print(f"OK 已处理 {n} 张图片 -> {outdir}")


if __name__ == "__main__":
    main()
""",
}

# ---------------- data-clean ----------------
SKILLS["data-clean"] = {
    "SKILL.md": """---
name: data-clean
description: 数据清洗：去重、去空、统一格式、修正类型、异常值处理，输出清洗报告
---

# 数据清洗

## 执行步骤
```bash
python "<仓库>\\data-clean\\scripts\\clean_data.py" --file 数据.xlsx --dedup --drop-empty --trim --output 清洗后.xlsx
```
- `--dedup`：按全部列去重；`--drop-empty`：删除整行全空；`--trim`：去除首尾空格。
- 脚本会打印清洗前后行数与各类修复计数，作为汇报数据。

## 注意事项
- 输出到工作区内；不修改原文件。
""",
    "scripts/clean_data.py": """#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# 数据清洗。
import argparse
import os
import sys

import pandas as pd


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--file", required=True)
    ap.add_argument("--output", default="")
    ap.add_argument("--dedup", action="store_true")
    ap.add_argument("--drop-empty", action="store_true")
    ap.add_argument("--trim", action="store_true")
    args = ap.parse_args()

    if args.file.lower().endswith(".csv"):
        df = pd.read_csv(args.file)
    else:
        df = pd.read_excel(args.file)
    before = len(df)
    report = []

    if args.trim:
        str_cols = df.select_dtypes(include="object").columns
        n = df[str_cols].apply(lambda s: s.astype(str).str.strip()).shape[0]
        for c in str_cols:
            df[c] = df[c].astype(str).str.strip()
        report.append(f"去空格: {n} 列")

    if args.drop_empty:
        n = df.dropna(how="all").shape[0]
        df = df.dropna(how="all")
        report.append(f"删整行空值: {before - n} 行")

    if args.dedup:
        n = df.drop_duplicates().shape[0]
        df = df.drop_duplicates()
        report.append(f"去重: {before - n} 行")

    print(f"清洗前: {before} 行 -> 清洗后: {len(df)} 行")
    for r in report:
        print("- " + r)
    if args.output:
        df.to_excel(args.output, index=False)
        print("OK 输出 -> " + args.output)
    else:
        print("OK (未指定 --output，仅打印)")


if __name__ == "__main__":
    main()
""",
}

# ---------------- archive-zip ----------------
SKILLS["archive-zip"] = {
    "SKILL.md": """---
name: archive-zip
description: 文件打包与解压：zip 压缩/解压，批量归档工作区文件
---

# 文件打包/解压

## 执行步骤
```bash
python "<仓库>\\archive-zip\\scripts\\zip_tools.py" zip --items 文件或目录... --output 归档.zip
python "<仓库>\\archive-zip\\scripts\\zip_tools.py" unzip --file 归档.zip --outdir 解压目录
python "<仓库>\\archive-zip\\scripts\\zip_tools.py" list --file 归档.zip
```
- 汇报压缩/解压的文件数与路径。

## 注意事项
- 只处理工作区内文件；解压目标必须在工作区内。
""",
    "scripts/zip_tools.py": """#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# zip 打包/解压/查看。
import argparse
import os
import sys
import zipfile


def main():
    ap = argparse.ArgumentParser()
    sub = ap.add_subparsers(dest="cmd", required=True)

    z = sub.add_parser("zip")
    z.add_argument("--items", nargs="*", required=True)
    z.add_argument("--output", required=True)

    u = sub.add_parser("unzip")
    u.add_argument("--file", required=True)
    u.add_argument("--outdir", default=".")

    l = sub.add_parser("list")
    l.add_argument("--file", required=True)

    args = ap.parse_args()

    if args.cmd == "zip":
        with zipfile.ZipFile(args.output, "w", zipfile.ZIP_DEFLATED) as zf:
            count = 0
            for item in args.items:
                if os.path.isdir(item):
                    for root, _, files in os.walk(item):
                        for f in files:
                            full = os.path.join(root, f)
                            zf.write(full, os.path.relpath(full, os.path.dirname(item) or "."))
                            count += 1
                else:
                    zf.write(item, os.path.basename(item))
                    count += 1
        print(f"OK 打包 {count} 个文件 -> {args.output}")

    elif args.cmd == "unzip":
        os.makedirs(args.outdir, exist_ok=True)
        with zipfile.ZipFile(args.file) as zf:
            zf.extractall(args.outdir)
        print(f"OK 解压 -> {args.outdir}")

    elif args.cmd == "list":
        with zipfile.ZipFile(args.file) as zf:
            for info in zf.infolist():
                print(f"{info.filename}  {info.file_size} B")


if __name__ == "__main__":
    main()
""",
}

# ---------------- python-script ----------------
SKILLS["python-script"] = {
    "SKILL.md": """---
name: python-script
description: 通用 Python 脚本编写与执行：把用户的一次性办公需求写成脚本并运行（工作区内）
---

# 通用 Python 脚本

## 适用场景
用户的需求不适合已有技能时，用本技能快速编写并运行 Python 脚本完成。

## 执行步骤
1. 拆解需求为步骤；选择内置库（pandas、openpyxl、python-docx、python-pptx、pypdf、PIL、matplotlib）。
2. 在工作区 `.oh_tmp` 目录编写脚本 `script_<任务>.py`（路径、输出明确）。
3. 运行：`python "<工作区>\\.oh_tmp\\script_xxx.py"`。
4. 校验输出存在；向用户汇报结果与关键数据；询问是否需要清理临时脚本。

## 注意事项
- 脚本只访问工作区内文件；删除/覆盖等破坏性操作先走审批。
- 输出用中文；出错时提炼「哪里错了 + 怎么改」，不输出大段堆栈。
- 任务完成后按需清理 `.oh_tmp`。
""",
}


def write_skill(name: str, files: dict):
    base = os.path.join(REPO, name)
    os.makedirs(base, exist_ok=True)
    for rel, content in files.items():
        path = os.path.join(base, rel)
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        print("  wrote", os.path.relpath(path, REPO))


def main():
    os.makedirs(REPO, exist_ok=True)
    print("安装技能到:", REPO)
    for name, files in SKILLS.items():
        print("-", name)
        write_skill(name, files)
    total = sum(len(v) for v in SKILLS.values())
    print(f"完成：{len(SKILLS)} 个技能，{total} 个文件")


if __name__ == "__main__":
    main()
